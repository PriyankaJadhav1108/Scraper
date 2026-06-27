"""
Convert Microsoft Office documents to PDF for hosted delivery.
"""

from __future__ import annotations

import base64
import html
import re
import shutil
import subprocess
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path


def find_libreoffice() -> str | None:
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/libreoffice",
        "/usr/bin/soffice",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def convert_with_libreoffice(source_path: Path, output_dir: Path) -> Path | None:
    soffice = find_libreoffice()
    if not soffice:
        return None

    output_dir.mkdir(parents=True, exist_ok=True)
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(source_path)],
            check=True,
            capture_output=True,
            timeout=180,
        )
    except Exception as exc:
        print(f"    ⚠  LibreOffice conversion failed for {source_path.name}: {exc}")
        return None

    pdf_path = output_dir / f"{source_path.stem}.pdf"
    return pdf_path if pdf_path.exists() else None


def _docx_to_html(source_path: Path) -> str:
    import mammoth

    with source_path.open("rb") as handle:
        result = mammoth.convert_to_html(handle)

    if result.messages:
        for message in result.messages[:3]:
            print(f"    ℹ  mammoth: {message}")

    return (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        "<style>body{font-family:Arial,sans-serif;margin:40px;line-height:1.5;}"
        "table{border-collapse:collapse;width:100%;}td,th{border:1px solid #ccc;padding:6px;}"
        "</style></head><body>"
        f"{result.value}</body></html>"
    )


def _pptx_to_html(source_path: Path) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(source_path))
    slides_html: list[str] = []

    for index, slide in enumerate(presentation.slides, 1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "text", "").strip():
                parts.append(f"<p>{html.escape(shape.text.strip())}</p>")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    encoded = base64.b64encode(image.blob).decode()
                    parts.append(
                        f'<img src="data:image/{image.ext};base64,{encoded}" '
                        'style="max-width:100%;margin:8px 0;" />'
                    )
                except Exception:
                    pass

        slides_html.append(
            f"<section class='slide'><h2>Slide {index}</h2>"
            f"{''.join(parts) or '<p></p>'}</section>"
        )

    return (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        "<style>body{font-family:Arial,sans-serif;margin:0;}"
        ".slide{page-break-after:always;padding:40px;min-height:90vh;}"
        "h2{color:#666;font-size:14px;}"
        "</style></head><body>"
        f"{''.join(slides_html)}</body></html>"
    )


def _html_to_pdf(html_content: str, pdf_path: Path) -> Path | None:
    from playwright.sync_api import sync_playwright

    def _render() -> Path | None:
        pdf_path.parent.mkdir(parents=True, exist_ok=True)
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1280, "height": 720})
            page.set_content(html_content, wait_until="networkidle")
            page.pdf(path=str(pdf_path), format="A4", print_background=True)
            browser.close()
        return pdf_path if pdf_path.exists() and pdf_path.stat().st_size > 0 else None

    with ThreadPoolExecutor(max_workers=1) as executor:
        return executor.submit(_render).result()


def convert_with_playwright(source_path: Path, output_dir: Path) -> Path | None:
    suffix = source_path.suffix.lower()
    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = output_dir / f"{source_path.stem}.pdf"

    try:
        if suffix == ".docx":
            html_content = _docx_to_html(source_path)
        elif suffix == ".pptx":
            html_content = _pptx_to_html(source_path)
        else:
            print(f"    ⚠  Unsupported Office format for Playwright conversion: {suffix}")
            return None

        with ThreadPoolExecutor(max_workers=1) as executor:
            return executor.submit(_html_to_pdf, html_content, pdf_path).result()
    except Exception as exc:
        print(f"    ⚠  Playwright conversion failed for {source_path.name}: {exc}")
        return None


def convert_html_file_to_pdf(source_path: Path, output_dir: Path) -> Path | None:
    raw = source_path.read_text(encoding="utf-8", errors="ignore")
    if "<TEXT>" in raw:
        start = raw.index("<TEXT>") + len("<TEXT>")
        end = raw.index("</TEXT>") if "</TEXT>" in raw else len(raw)
        raw = raw[start:end]
    body_match = re.search(r"<body[^>]*>(.*)</body>", raw, re.I | re.S)
    html_content = body_match.group(1) if body_match else raw
    wrapped = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        "<style>body{font-family:'Times New Roman',serif;margin:40px;line-height:1.4;}"
        "table{border-collapse:collapse;width:100%;}td,th{border:1px solid #ccc;padding:4px;}"
        "</style></head><body>"
        f"{html_content}</body></html>"
    )
    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = output_dir / f"{source_path.stem}.pdf"
    return _html_to_pdf(wrapped, pdf_path)


def convert_office_to_pdf(source_path: Path, output_dir: Path) -> Path | None:
    if source_path.suffix.lower() == ".pdf":
        return source_path
    if source_path.suffix.lower() in {".htm", ".html"}:
        return convert_html_file_to_pdf(source_path, output_dir)

    pdf_path = convert_with_libreoffice(source_path, output_dir)
    if pdf_path:
        return pdf_path

    print(f"    ℹ  LibreOffice unavailable; converting {source_path.name} via Playwright")
    return convert_with_playwright(source_path, output_dir)
